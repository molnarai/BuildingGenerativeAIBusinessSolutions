#!/usr/bin/env python
TITLE=r"""
  ____  ____  _____   _          ____  ____ _______  __   ____                          _            
 |  _ \|  _ \|  ___| | |_ ___   |  _ \|  _ \_   _\ \/ /  / ___|___  _ ____   _____ _ __| |_ ___ _ __ 
 | |_) | | | | |_    | __/ _ \  | |_) | |_) || |  \  /  | |   / _ \| '_ \ \ / / _ \ '__| __/ _ \ '__|
 |  __/| |_| |  _|   | || (_) | |  __/|  __/ | |  /  \  | |__| (_) | | | \ V /  __/ |  | ||  __/ |   
 |_|   |____/|_|      \__\___/  |_|   |_|    |_| /_/\_\  \____\___/|_| |_|\_/ \___|_|   \__\___|_|   
                                                                                                     
"""
import os
import sys
from pathlib import Path
import argparse
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE

def main():
    print(TITLE)
    parser = argparse.ArgumentParser()
    parser.add_argument('-f', '--file', help='pdf file to convert', required=True)
    parser.add_argument('--template', help='pptx template file to use', required=False)
    parser.add_argument('--format', help='slide format: 16x9 or 4x3 (default: 16x9)', 
                        choices=['16x9', '4x3'], default='16x9')
    parser.add_argument('--density', help='DPI for PNG conversion (default: 200)', 
                        type=int, default=200)
    parser.add_argument('--clear', help='Remove template slide elements before adding background', 
                        action='store_true')
    args = parser.parse_args()
    
    input_file = Path(args.file)
    if not input_file.exists():
        print("File does not exist")
        sys.exit(1)
    
    if input_file.suffix.lower() != '.pdf':
        print("Input file must be a PDF")
        sys.exit(1)
    
    # Convert PDF pages to PNG images with density 200
    print(f"Converting PDF pages to PNG images (DPI: {args.density})...")
    images = convert_from_path(str(input_file), dpi=args.density)
    
    # Create or load presentation
    if args.template and Path(args.template).exists():
        prs = Presentation(args.template)
        print(f"Using template: {args.template}")
        # Get the last slide's layout to use as template
        if len(prs.slides) > 0:
            template_slide = prs.slides[-1]
            template_layout = template_slide.slide_layout
            use_template_slide = True
            print(f"Using last slide as template (layout: {template_layout.name})")
        else:
            template_layout = prs.slide_layouts[6]  # Blank layout fallback
            use_template_slide = False
            print("Template has no slides, using blank layout")
    else:
        prs = Presentation()
        # Set slide dimensions based on format
        if args.format == '16x9':
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
        else:  # 4x3
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
        print(f"Using {args.format} format")
        template_layout = prs.slide_layouts[6]  # Blank layout
        use_template_slide = False
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Create slides with PNG backgrounds
    print(f"Creating PPTX with {len(images)} slides...")
    for i, image in enumerate(images):
        # Save PNG temporarily
        temp_png = f"temp_page_{i}.png"
        image.save(temp_png, "PNG")
        
        # Add slide using template layout
        slide = prs.slides.add_slide(template_layout)
        
        # Copy shapes from template slide if using template and not clearing
        if use_template_slide and not args.clear:
            for shape in template_slide.shapes:
                # Handle picture shapes specially to preserve images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Get the image from the original shape
                    rId = shape._element.blipFill.blip.rEmbed
                    image_part = template_slide.part.related_part(rId)
                    # Add the same image to the new slide using image bytes
                    from io import BytesIO
                    pic = slide.shapes.add_picture(
                        BytesIO(image_part.blob),
                        shape.left, shape.top,
                        shape.width, shape.height
                    )
                else:
                    # Copy other shape types normally
                    el = shape.element
                    newel = deepcopy(el)
                    slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        # Clear template elements if --clear option is set
        if args.clear:
            shapes_to_remove = list(slide.shapes)
            for shape in shapes_to_remove:
                sp = shape.element
                sp.getparent().remove(sp)
        
        # Insert picture as background
        picture = slide.shapes.add_picture(temp_png, 0, 0, width=slide_width, height=slide_height)
        # Move picture to back (behind all other shapes)
        slide.shapes._spTree.remove(picture._element)
        slide.shapes._spTree.insert(2, picture._element)
        
        # Clean up temporary PNG
        os.remove(temp_png)
        
        print(f"  Added slide {i+1}/{len(images)}")
    
    # Save PPTX with same name as input file
    output_file = input_file.with_suffix('.pptx')
    prs.save(str(output_file))
    print(f"\nSuccessfully created: {output_file}")

if __name__ == "__main__":
    main()

