#!/usr/bin/env python3
"""
Document content extraction module.
Handles extraction from various document formats.
"""

import logging
from pathlib import Path
from typing import Dict, Any, Optional
import mimetypes

logger = logging.getLogger(__name__)

class DocumentExtractor:
    """Extracts content from various document types."""

    def __init__(self):
        self.extractors = {
            'application/pdf': self._extract_pdf,
            'image/jpeg': self._extract_image_ocr,
            'image/png': self._extract_image_ocr,
            'image/tiff': self._extract_image_ocr,
            'text/plain': self._extract_text,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
        }

    def extract_content(self, file_path: Path) -> Dict[str, Any]:
        """
        Extract content from a document file.

        Args:
            file_path: Path to the document

        Returns:
            Dictionary containing extracted content and metadata
        """
        # Detect MIME type
        mime_type, _ = mimetypes.guess_type(str(file_path))

        if not mime_type or mime_type not in self.extractors:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")

        # Extract content
        extractor_func = self.extractors[mime_type]
        content = extractor_func(file_path)

        return {
            'file_name': file_path.name,
            'file_path': str(file_path),
            'mime_type': mime_type,
            'content': content['text'],
            'metadata': content.get('metadata', {}),
            'extraction_method': content.get('method', 'unknown')
        }

    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PDF using PyMuPDF (fitz)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            raise ImportError("PyMuPDF is required for PDF extraction. Install with: pip install PyMuPDF")

        doc = fitz.open(str(file_path))
        text_parts = []
        metadata = {}

        # Extract metadata
        if doc.metadata:
            metadata = {
                'title': doc.metadata.get('title', ''),
                'author': doc.metadata.get('author', ''),
                'subject': doc.metadata.get('subject', ''),
                'creator': doc.metadata.get('creator', ''),
                'producer': doc.metadata.get('producer', ''),
                'creation_date': doc.metadata.get('creationDate', ''),
                'modification_date': doc.metadata.get('modDate', ''),
                'page_count': len(doc)
            }

        # Extract text from each page
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text()
            text_parts.append(f"--- Page {page_num + 1} ---\n{text}")

        doc.close()

        return {
            'text': '\n\n'.join(text_parts),
            'metadata': metadata,
            'method': 'pdf_text_extraction'
        }

    def _extract_image_ocr(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from images using OCR."""
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            raise ImportError("Pillow and pytesseract are required for OCR. Install with: pip install Pillow pytesseract")

        image = Image.open(str(file_path))
        text = pytesseract.image_to_string(image)

        metadata = {
            'image_size': image.size,
            'image_mode': image.mode,
            'format': image.format
        }

        return {
            'text': text,
            'metadata': metadata,
            'method': 'ocr'
        }

    def _extract_text(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            # Try with different encodings
            for encoding in ['latin-1', 'cp1252', 'utf-16']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError(f"Could not decode text file: {file_path}")

        file_size = file_path.stat().st_size

        return {
            'text': text,
            'metadata': {'file_size': file_size},
            'method': 'direct_text'
        }

    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from DOCX files."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX extraction. Install with: pip install python-docx")

        doc = Document(str(file_path))
        text_parts = []

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        metadata = {
            'paragraph_count': len(text_parts),
            'core_properties': {
                'title': doc.core_properties.title or '',
                'author': doc.core_properties.author or '',
                'subject': doc.core_properties.subject or '',
                'created': str(doc.core_properties.created or ''),
                'modified': str(doc.core_properties.modified or '')
            }
        }

        return {
            'text': '\n\n'.join(text_parts),
            'metadata': metadata,
            'method': 'docx_extraction'
        }

    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX extraction. Install with: pip install python-pptx")

        prs = Presentation(str(file_path))
        text_parts = []
        slide_count = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
                slide_count += 1

        metadata = {
            'slide_count': slide_count,
            'total_slides': len(prs.slides)
        }

        return {
            'text': '\n\n'.join(text_parts),
            'metadata': metadata,
            'method': 'pptx_extraction'
        }